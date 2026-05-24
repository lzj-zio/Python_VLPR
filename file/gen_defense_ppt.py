"""
答辩PPT生成脚本 — 基于实际项目产物
"基于OpenCV的工程车车牌识别系统设计与实现"
Author: 李志杰 | 2022b15080 | 物联网工程 | 浙江水利水电学院

本脚本以实际项目代码实现为基准生成答辩PPT，确保与演示产物一致。
核心实现：OpenCV图像处理 + HOG特征 + SVM分类 + Tkinter GUI
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# =========================================================
# COLOR SYSTEM
# =========================================================
BG      = RGBColor(0x0D, 0x1B, 0x3E)
BG2     = RGBColor(0x13, 0x25, 0x52)
CARD    = RGBColor(0x18, 0x2D, 0x60)
CARD2   = RGBColor(0x0F, 0x2B, 0x4E)
BLUE    = RGBColor(0x4A, 0x9E, 0xFF)
CYAN    = RGBColor(0x00, 0xD4, 0xC8)
GREEN   = RGBColor(0x39, 0xD3, 0x7A)
ORANGE  = RGBColor(0xFF, 0x8C, 0x42)
PURPLE  = RGBColor(0xA0, 0x6C, 0xFF)
RED     = RGBColor(0xFF, 0x4D, 0x6D)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0xB0, 0xBE, 0xD9)
DGRAY   = RGBColor(0x5E, 0x72, 0x9A)
GOLD    = RGBColor(0xFF, 0xD7, 0x00)
LGOLD   = RGBColor(0xFF, 0xE8, 0x80)

# =========================================================
# HELPERS
# =========================================================
def set_bg(slide, color):
    bg = slide.background
    f = bg.fill
    f.solid()
    f.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line_color=None, line_w=1.5, shape=MSO_SHAPE.RECTANGLE, adj=None):
    s = slide.shapes.add_shape(shape, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    if adj is not None:
        try:
            s.adjustments[0] = adj
        except:
            pass
    return s

def rbox(slide, l, t, w, h, fill=None, line_color=None, line_w=1.5):
    return box(slide, l, t, w, h, fill, line_color, line_w, MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.04)

def txt(slide, l, t, w, h, text, size=14, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, wrap=True, italic=False, font='Microsoft YaHei'):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return tb

def circle(slide, l, t, size, fill, text='', fsize=13):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if text:
        tf = s.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(fsize)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = 'Microsoft YaHei'
    return s

def divider(slide, x, y, w, color=BLUE, thick=2):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(thick))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def arrow_right(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def arrow_down(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def tag(slide, l, t, label, fill, text_color=WHITE, fsize=11):
    w = Inches(0.9)
    h = Inches(0.28)
    b = rbox(slide, l, t, w, h, fill=fill)
    txt(slide, l, t, w, h, label, size=fsize, color=text_color, bold=True, align=PP_ALIGN.CENTER)
    return b

def section_title(slide, title, subtitle=''):
    set_bg(slide, BG)
    box(slide, Inches(0), Inches(0), Inches(0.07), Inches(7.5), fill=BLUE)
    txt(slide, Inches(0.3), Inches(0.28), Inches(10), Inches(0.6),
        title, size=26, color=WHITE, bold=True)
    if subtitle:
        txt(slide, Inches(0.3), Inches(0.88), Inches(11), Inches(0.4),
            subtitle, size=14, color=GRAY)
    divider(slide, Inches(0.3), Inches(1.3), Inches(12.7), color=BLUE, thick=2)

# =========================================================
# SLIDE 1 — COVER
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG)

box(slide, 0, 0, prs.slide_width, Inches(0.08), fill=BLUE)
box(slide, 0, Inches(0.08), prs.slide_width, Inches(0.04), fill=CYAN)

big_circle = box(slide, Inches(8.5), Inches(1.0), Inches(4.8), Inches(4.8),
                 fill=RGBColor(0x14, 0x2D, 0x6E), shape=MSO_SHAPE.OVAL)
big_circle.line.color.rgb = BLUE
big_circle.line.width = Pt(1)
box(slide, Inches(9.2), Inches(1.7), Inches(3.4), Inches(3.4),
    fill=RGBColor(0x1A, 0x38, 0x80), shape=MSO_SHAPE.OVAL).line.fill.background()

# 实际技术栈
txt(slide, Inches(9.0), Inches(2.4), Inches(3.8), Inches(1.8),
    'OpenCV\n+\nSVM', size=17, color=CYAN, bold=True,
    align=PP_ALIGN.CENTER)

box(slide, Inches(0.5), Inches(0.5), Inches(0.05), Inches(0.8), fill=BLUE)

txt(slide, Inches(0.7), Inches(0.85), Inches(8), Inches(1.1),
    '基于OpenCV的工程车', size=36, color=WHITE, bold=True)
txt(slide, Inches(0.7), Inches(1.9), Inches(8), Inches(0.9),
    '车牌识别系统设计与实现', size=36, color=WHITE, bold=True)

divider(slide, Inches(0.7), Inches(2.85), Inches(7.5), color=CYAN, thick=3)

txt(slide, Inches(0.7), Inches(3.0), Inches(8), Inches(0.5),
    'Design and Implementation of Engineering Vehicle License Plate Recognition System Based on OpenCV',
    size=11, color=DGRAY, italic=True)

info_y = Inches(3.8)
rbox(slide, Inches(0.7), info_y, Inches(6.8), Inches(2.8), fill=CARD, line_color=BLUE, line_w=1)

fields = [
    ('学    院', '计算机科学与技术'),
    ('专    业', '物联网工程'),
    ('学    号', '2022b15080'),
    ('学生姓名', '李志杰'),
    ('指导教师', '徐欧官'),
]
for i, (k, v) in enumerate(fields):
    fy = info_y + Inches(0.12) + i * Inches(0.5)
    txt(slide, Inches(1.0), fy, Inches(1.6), Inches(0.42),
        k, size=13, color=GRAY)
    txt(slide, Inches(2.8), fy, Inches(4.0), Inches(0.42),
        v, size=13, color=WHITE, bold=True)
    if i < len(fields)-1:
        divider(slide, Inches(1.0), fy + Inches(0.42), Inches(5.8), color=CARD2, thick=1)

txt(slide, Inches(0.7), Inches(6.75), Inches(6), Inches(0.45),
    '浙江水利水电学院  ·  2026届', size=13, color=DGRAY)

# =========================================================
# SLIDE 2 — OUTLINE
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title(slide, '目  录', 'Table of Contents')

items = [
    ('01', '研究背景与意义', '工程车识别挑战 · 行业需求 · 技术选型', BLUE),
    ('02', '系统相关技术', 'OpenCV · HOG特征 · SVM分类 · 图像处理', CYAN),
    ('03', '系统总体设计', '双通道架构 · 模块划分 · 多平台部署', GREEN),
    ('04', '系统详细实现', '预处理 · 定位 · 识别 · GUI界面', ORANGE),
    ('05', '系统测试与分析', '分模块测试 · 摄像头模式 · Web仪表盘', PURPLE),
    ('06', '总结与展望', '成果归纳 · 局限分析 · 优化方向', GOLD),
]

cols = 3
for i, (num, title, sub, color) in enumerate(items):
    col = i % cols
    row = i // cols
    cx = Inches(0.5) + col * Inches(4.2)
    cy = Inches(1.6) + row * Inches(2.5)
    rbox(slide, cx, cy, Inches(4.0), Inches(2.1), fill=CARD, line_color=color, line_w=1.5)
    circle(slide, cx + Inches(0.18), cy + Inches(0.18), Inches(0.55), color, num, fsize=18)
    txt(slide, cx + Inches(0.85), cy + Inches(0.2), Inches(3.0), Inches(0.45),
        title, size=16, color=WHITE, bold=True)
    txt(slide, cx + Inches(0.85), cy + Inches(0.68), Inches(2.9), Inches(0.5),
        sub, size=11, color=GRAY)
    divider(slide, cx + Inches(0.18), cy + Inches(1.25), Inches(3.6), color=color, thick=1)

# =========================================================
# SLIDE 3 — BACKGROUND & MOTIVATION
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title(slide, '01  研究背景与意义', '工程车在复杂工地场景下的车牌识别挑战')

# 核心痛点
rbox(slide, Inches(0.3), Inches(1.55), Inches(5.9), Inches(5.4), fill=CARD, line_color=BLUE, line_w=1)
txt(slide, Inches(0.5), Inches(1.7), Inches(4), Inches(0.4),
    '核心痛点', size=16, color=BLUE, bold=True)
divider(slide, Inches(0.5), Inches(2.12), Inches(5.5), color=BLUE, thick=1)

problems = [
    ('图像模糊', '运动抖动 / 低分辨率 / 高速行驶'),
    ('光照不均', '强逆光 / 夜间 / 雨雾遮挡场景'),
    ('车牌倾斜', '大角度停放 / 工地起伏路面'),
    ('字符污损', '泥水覆盖 / 标签脱落 / 划痕腐蚀'),
    ('多类型车牌', '蓝牌 / 黄牌 / 新能源绿牌处理'),
]
for i, (prob, desc) in enumerate(problems):
    py = Inches(2.3) + i * Inches(0.68)
    circle(slide, Inches(0.55), py + Pt(3), Inches(0.32), RED, str(i+1), fsize=11)
    txt(slide, Inches(1.0), py, Inches(2.0), Inches(0.35),
        prob, size=13, color=ORANGE, bold=True)
    txt(slide, Inches(3.1), py, Inches(2.8), Inches(0.35),
        desc, size=12, color=GRAY)

# 行业背景
rbox(slide, Inches(6.5), Inches(1.55), Inches(6.5), Inches(2.5), fill=CARD, line_color=CYAN, line_w=1)
txt(slide, Inches(6.7), Inches(1.7), Inches(5), Inches(0.4),
    '行业需求背景', size=16, color=CYAN, bold=True)
divider(slide, Inches(6.7), Inches(2.12), Inches(6.1), color=CYAN, thick=1)

stats = [
    ('950万+', '全国工程机械保有量（2025）', BLUE),
    ('65%', '大型施工项目计划部署智能识别', GREEN),
    ('6.8%', '年均保有量增长率', ORANGE),
]
for i, (num, label, color) in enumerate(stats):
    sx = Inches(6.7) + i * Inches(2.1)
    txt(slide, sx, Inches(2.3), Inches(2.0), Inches(0.65),
        num, size=26, color=color, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, sx, Inches(2.95), Inches(2.0), Inches(0.8),
        label, size=10, color=GRAY, align=PP_ALIGN.CENTER)

# 技术选型依据
rbox(slide, Inches(6.5), Inches(4.3), Inches(6.5), Inches(2.6), fill=CARD, line_color=GREEN, line_w=1)
txt(slide, Inches(6.7), Inches(4.45), Inches(5), Inches(0.4),
    '技术选型依据', size=16, color=GREEN, bold=True)
divider(slide, Inches(6.7), Inches(4.87), Inches(6.1), color=GREEN, thick=1)

values = [
    'OpenCV图像处理：成熟稳定，无需GPU，适合嵌入式部署',
    'HOG+SVM分类：训练数据需求少，推理速度快',
    '双通道互补定位：形状+颜色双重验证，提升鲁棒性',
    '轻量化架构：适配桌面端、手机端、Web仪表盘',
]
for i, v in enumerate(values):
    circle(slide, Inches(6.7), Inches(5.05) + i * Inches(0.47), Inches(0.28), GREEN, '✓', fsize=10)
    txt(slide, Inches(7.1), Inches(5.05) + i * Inches(0.47), Inches(5.6), Inches(0.42),
        v, size=12, color=GRAY)

# =========================================================
# SLIDE 4 — TECH STACK OVERVIEW
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title(slide, '02  核心技术栈', 'OpenCV图像处理 + HOG特征提取 + SVM分类识别')

# 技术栈全景 — 基于实际代码
tech_layers = [
    ('图像采集层', ['Tkinter GUI (桌面端)', 'Kivy GUI (Android端)', 'USB摄像头实时采集'], BLUE),
    ('图像处理层', ['高斯模糊降噪', 'HSV颜色空间转换', '形态学开/闭操作', 'Canny边缘检测'], CYAN),
    ('车牌定位层', ['形状通道: 灰度→边缘→轮廓筛选', '颜色通道: HSV三色掩码提取', '仿射变换倾斜校正', '投影直方图字符分割'], GREEN),
    ('字符识别层', ['HOG 64维特征提取', 'RBF-SVM 中文省份分类(31类)', 'RBF-SVM 字母数字分类(34类)', '多帧投票融合机制'], ORANGE),
    ('数据展示层', ['Tkinter双通道结果对比', 'SQLite识别记录持久化', 'WebSocket实时推送', 'Chart.js可视化仪表盘'], PURPLE),
]

lw = Inches(12.5)
lx = Inches(0.4)

for i, (name, items, color) in enumerate(tech_layers):
    ly = Inches(1.5) + i * Inches(1.15)
    rbox(slide, lx, ly, lw, Inches(1.02), fill=CARD, line_color=color, line_w=1.5)
    rbox(slide, lx + Inches(0.12), ly + Inches(0.22), Inches(1.6), Inches(0.5), fill=color)
    txt(slide, lx + Inches(0.12), ly + Inches(0.22), Inches(1.6), Inches(0.5),
        name, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 用分隔符展示各技术项
    items_text = '  |  '.join(items)
    txt(slide, lx + Inches(1.9), ly + Inches(0.28), Inches(10.3), Inches(0.6),
        items_text, size=11, color=GRAY)
    if i < len(tech_layers) - 1:
        arrow_down(slide, Inches(5.8), ly + Inches(1.02), Inches(0.25), Inches(0.12), DGRAY)

# =========================================================
# SLIDE 5 — SYSTEM ARCHITECTURE
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title(slide, '03  系统总体架构', '双通道互补定位 · SVM分类识别 · 多平台部署')

# 实际架构：双通道定位 + SVM识别
layers = [
    ('图像输入模块', '本地图片选择 / USB摄像头实时采集',
     'Tkinter文件对话框 / cv2.VideoCapture (2秒/帧)', CYAN),
    ('双通道车牌定位', '形状通道 + 颜色通道 并行执行，结果互为验证',
     '轮廓筛选(宽高比2~5.5) + HSV蓝/黄/绿三色掩码', BLUE),
    ('字符分割与识别', '仿射变换校正 → 垂直投影分割 → HOG特征 → SVM预测',
     '64维HOG特征 + RBF-SVM (C=1, γ=0.5)', GREEN),
    ('结果输出模块', 'GUI双通道对比展示 / SQLite记录 / WebSocket推送',
     '多帧融合(窗口=5, 阈值=3) / 实时仪表盘', ORANGE),
]

lw = Inches(12.5)
lh = Inches(1.3)
lx = Inches(0.4)

for i, (name, content, tech, color) in enumerate(layers):
    ly = Inches(1.55) + i * Inches(1.5)
    rbox(slide, lx, ly, lw, lh, fill=CARD, line_color=color, line_w=2)
    rbox(slide, lx + Inches(0.15), ly + Inches(0.35), Inches(2.0), Inches(0.5),
         fill=color)
    txt(slide, lx + Inches(0.15), ly + Inches(0.35), Inches(2.0), Inches(0.5),
        name, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, lx + Inches(2.4), ly + Inches(0.15), Inches(6.5), Inches(0.5),
        content, size=13, color=WHITE, bold=True)
    txt(slide, lx + Inches(2.4), ly + Inches(0.7), Inches(6.5), Inches(0.5),
        tech, size=10, color=GRAY)
    if i < 3:
        arrow_down(slide, Inches(5.8), ly + lh + Inches(0.08),
                   Inches(0.25), Inches(0.2), DGRAY)

# 多平台部署
mod_y = Inches(7.9) - Inches(1.8)
txt(slide, Inches(0.4), mod_y, Inches(12), Inches(0.38),
    '多平台部署方案', size=14, color=DGRAY, bold=True)

modules = [
    ('桌面端 (Tkinter)', '双通道对比 · 摄像头模式 · 记录查询', BLUE),
    ('Android端 (Kivy)', 'buildozer打包 · 骁龙865 · 实时识别', GREEN),
    ('Web仪表盘', 'Chart.js可视化 · WebSocket推送 · 历史查询', PURPLE),
]
mw = Inches(3.8)
for i, (label, sub, color) in enumerate(modules):
    mx = Inches(0.4) + i * Inches(4.1)
    rbox(slide, mx, mod_y + Inches(0.42), mw, Inches(0.78), fill=CARD2, line_color=color, line_w=1)
    txt(slide, mx, mod_y + Inches(0.44), mw, Inches(0.32),
        label, size=12, color=color, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, mx, mod_y + Inches(0.76), mw, Inches(0.28),
        sub, size=10, color=DGRAY, align=PP_ALIGN.CENTER, italic=True)

# =========================================================
# SLIDE 6 — IMAGE PREPROCESSING (实际代码实现)
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title(slide, '04  图像预处理模块', '基于实际代码的预处理流程：高斯模糊 → 灰度化 → 形态学操作 → 边缘检测')

# 实际预处理流程
steps = [
    ('原图输入\n缩放至≤1000px', DGRAY),
    ('高斯模糊\nkernel=5', BLUE),
    ('灰度化 +\n形态学开操作', CYAN),
    ('加权混合\naddWeighted', GREEN),
    ('Otsu二值化\n阈值分割', ORANGE),
    ('Canny边缘\n100/200', PURPLE),
]
sw = Inches(1.7)
sy = Inches(1.6)
sh = Inches(1.3)
sx0 = Inches(0.35)

for i, (label, color) in enumerate(steps):
    sx = sx0 + i * (sw + Inches(0.5))
    rbox(slide, sx, sy, sw, sh, fill=color)
    txt(slide, sx, sy, sw, sh, label, size=11, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        arrow_right(slide, sx + sw, sy + Inches(0.45), Inches(0.45), Inches(0.35), DGRAY)

# 左侧：预处理细节（实际代码参数）
rbox(slide, Inches(0.3), Inches(3.15), Inches(5.8), Inches(4.0), fill=CARD, line_color=BLUE, line_w=1.5)
txt(slide, Inches(0.5), Inches(3.3), Inches(4), Inches(0.4),
    '预处理参数配置', size=13, color=BLUE, bold=True)
divider(slide, Inches(0.5), Inches(3.72), Inches(5.4), color=BLUE, thick=1)

preprocess_pts = [
    '图像缩放：最大宽度限制1000px，保持宽高比',
    '高斯模糊：5×5高斯核，抑制高频噪声',
    '形态学开操作：20×20核，去除细小噪点',
    '加权混合：addWeighted融合模糊与原图，保留边缘',
    'Otsu二值化：自动阈值，适应不同光照条件',
    'Canny边缘检测：低阈值100、高阈值200',
]
for i, pt in enumerate(preprocess_pts):
    circle(slide, Inches(0.5), Inches(3.9) + i * Inches(0.52), Inches(0.26), BLUE, str(i+1), fsize=9)
    txt(slide, Inches(0.88), Inches(3.9) + i * Inches(0.52), Inches(5.0), Inches(0.48),
        pt, size=11, color=GRAY)

# 右侧：形态学闭+开操作
rbox(slide, Inches(6.6), Inches(3.15), Inches(6.4), Inches(4.0), fill=CARD, line_color=CYAN, line_w=1.5)
txt(slide, Inches(6.8), Inches(3.3), Inches(5.5), Inches(0.4),
    '形态学后处理', size=13, color=CYAN, bold=True)
divider(slide, Inches(6.8), Inches(3.72), Inches(6.0), color=CYAN, thick=1)

morph_pts = [
    '闭操作：4×19水平核，连接断裂的字符笔画',
    '开操作：去除细小干扰区域',
    '轮廓查找：面积阈值>2000px，宽高比2~5.5',
    '最小外接矩形：获取精确车牌边界框',
    '仿射变换：基于boxPoints四点校正倾斜',
    'HSV颜色分类：自动判断蓝/黄/绿牌类型',
]
for i, pt in enumerate(morph_pts):
    circle(slide, Inches(6.8), Inches(3.9) + i * Inches(0.52), Inches(0.26), CYAN, str(i+1), fsize=9)
    txt(slide, Inches(7.18), Inches(3.9) + i * Inches(0.52), Inches(5.5), Inches(0.48),
        pt, size=11, color=GRAY)

# =========================================================
# SLIDE 7 — DUAL-CHANNEL PLATE LOCALIZATION (实际实现)
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title(slide, '04  双通道车牌定位', '形状通道 + 颜色通道 并行执行，结果互为验证')

# 左侧：形状定位通道
rbox(slide, Inches(0.3), Inches(1.55), Inches(6.2), Inches(5.4), fill=CARD, line_color=BLUE, line_w=1.5)
txt(slide, Inches(0.5), Inches(1.7), Inches(5), Inches(0.4),
    '形状定位通道（img_color_contours）', size=14, color=BLUE, bold=True)
divider(slide, Inches(0.5), Inches(2.12), Inches(5.8), color=BLUE, thick=1)

loc_steps = [
    ('步骤1', '高斯模糊 → 灰度化 → 形态学开操作(20×20)', BLUE),
    ('步骤2', '加权混合原始图像，增强边缘保留', CYAN),
    ('步骤3', 'Otsu二值化 → Canny边缘检测(100/200)', GREEN),
    ('步骤4', '形态学闭+开操作(4×19水平核)', ORANGE),
    ('步骤5', '轮廓筛选：面积>2000 & 宽高比2~5.5', PURPLE),
    ('步骤6', '仿射变换校正倾斜 → HSV颜色分类', RED),
]
for i, (step, desc, color) in enumerate(loc_steps):
    ny = Inches(2.28) + i * Inches(0.6)
    rbox(slide, Inches(0.5), ny, Inches(5.8), Inches(0.52), fill=CARD2, line_color=color, line_w=1)
    tag(slide, Inches(0.55), ny + Inches(0.1), step, color, fsize=9)
    txt(slide, Inches(1.55), ny + Inches(0.08), Inches(4.5), Inches(0.35),
        desc, size=10.5, color=GRAY)

# 右侧：颜色定位通道
rbox(slide, Inches(6.8), Inches(1.55), Inches(6.2), Inches(5.4), fill=CARD, line_color=CYAN, line_w=1.5)
txt(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.4),
    '颜色定位通道（img_only_color）', size=14, color=CYAN, bold=True)
divider(slide, Inches(7.0), Inches(2.12), Inches(5.8), color=CYAN, thick=1)

# HSV颜色阈值
color_ranges = [
    ('蓝色车牌', 'H[100~130]  S[110~255]  V[110~255]', BLUE),
    ('黄色车牌', 'H[15~50]    S[55~255]   V[55~255]', ORANGE),
    ('绿色车牌', 'H[50~100]   S[50~255]   V[50~255]', GREEN),
]
for i, (name, range_val, color) in enumerate(color_ranges):
    ry = Inches(2.3) + i * Inches(0.72)
    rbox(slide, Inches(7.0), ry, Inches(5.8), Inches(0.6), fill=CARD2, line_color=color, line_w=1)
    tag(slide, Inches(7.05), ry + Inches(0.14), name, color, fsize=9)
    txt(slide, Inches(8.05), ry + Inches(0.12), Inches(4.5), Inches(0.35),
        range_val, size=11, color=GRAY)

# 后续处理
txt(slide, Inches(7.0), Inches(4.5), Inches(5.5), Inches(0.4),
    '后续处理流程', size=13, color=CYAN, bold=True)
divider(slide, Inches(7.0), Inches(4.9), Inches(5.8), color=CYAN, thick=1)

post_steps = [
    '三色掩码合并 → 形态学闭+开操作去噪',
    '轮廓查找 → 宽高比/面积几何约束筛选',
    '水平投影直方图定位车牌行区域',
    '垂直投影直方图分割7个字符',
    '判尾过滤：去除过细的误识别"1"',
]
for i, pt in enumerate(post_steps):
    circle(slide, Inches(7.0), Inches(5.05) + i * Inches(0.42), Inches(0.26), CYAN, str(i+1), fsize=9)
    txt(slide, Inches(7.38), Inches(5.05) + i * Inches(0.42), Inches(5.3), Inches(0.38),
        pt, size=10.5, color=GRAY)

# =========================================================
# SLIDE 8 — CHARACTER RECOGNITION (实际SVM+HOG)
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title(slide, '04  字符识别模块', 'HOG 64维特征提取 + RBF-SVM分类 · 双模型架构（中文省份 + 字母数字）')

# 左侧：HOG特征提取
rbox(slide, Inches(0.3), Inches(1.55), Inches(5.9), Inches(5.4), fill=CARD, line_color=GREEN, line_w=1.5)
txt(slide, Inches(0.5), Inches(1.7), Inches(5.0), Inches(0.4),
    'HOG特征提取流程', size=14, color=GREEN, bold=True)
divider(slide, Inches(0.5), Inches(2.12), Inches(5.5), color=GREEN, thick=1)

hog_stages = [
    ('Step 1', '字符归一化', '字符图像补齐 → resize至20×20像素', BLUE),
    ('Step 2', '矩校正(deskew)', 'cv2.moments计算倾斜角 → warpAffine反向映射校正', CYAN),
    ('Step 3', 'Sobel梯度计算', 'Sobel算子提取水平/垂直梯度幅值与方向', ORANGE),
    ('Step 4', 'HOG直方图统计', '16方向bin × 4个cell(10×10) = 64维特征向量', GREEN),
    ('Step 5', 'Hellinger归一化', 'sqrt归一化 → L2归一化，消除光照影响', PURPLE),
]
for i, (step, name, desc, color) in enumerate(hog_stages):
    ny = Inches(2.28) + i * Inches(0.9)
    rbox(slide, Inches(0.5), ny, Inches(5.5), Inches(0.78), fill=CARD2, line_color=color, line_w=1)
    tag(slide, Inches(0.55), ny + Inches(0.22), step, color, fsize=9)
    txt(slide, Inches(1.55), ny + Inches(0.08), Inches(4.3), Inches(0.28),
        name, size=12, color=color, bold=True)
    txt(slide, Inches(1.55), ny + Inches(0.4), Inches(4.3), Inches(0.35),
        desc, size=10, color=GRAY)

# 右侧：SVM模型
rbox(slide, Inches(6.6), Inches(1.55), Inches(6.4), Inches(2.5), fill=CARD, line_color=BLUE, line_w=1.5)
txt(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(0.4),
    'SVM双模型架构', size=14, color=BLUE, bold=True)
divider(slide, Inches(6.8), Inches(2.12), Inches(6.0), color=BLUE, thick=1)

# 两个模型信息
models_info = [
    ('中文省份模型', 'svmchinese.dat (3.5MB)', '31类省份简称 · 3,232张训练样本', ORANGE),
    ('字母数字模型', 'svm.dat (4.5MB)', '34类(0-9,A-Z) · 13,163张训练样本', CYAN),
]
for i, (name, file_info, detail, color) in enumerate(models_info):
    my = Inches(2.3) + i * Inches(0.75)
    rbox(slide, Inches(6.8), my, Inches(6.0), Inches(0.65), fill=CARD2, line_color=color, line_w=1)
    tag(slide, Inches(6.85), my + Inches(0.18), name, color, fsize=9)
    txt(slide, Inches(7.85), my + Inches(0.05), Inches(4.8), Inches(0.28),
        file_info, size=11, color=WHITE, bold=True)
    txt(slide, Inches(7.85), my + Inches(0.35), Inches(4.8), Inches(0.25),
        detail, size=10, color=GRAY)

# SVM参数
rbox(slide, Inches(6.6), Inches(4.3), Inches(6.4), Inches(2.65), fill=CARD, line_color=CYAN, line_w=1.5)
txt(slide, Inches(6.8), Inches(4.45), Inches(5), Inches(0.4),
    'SVM参数与训练配置', size=14, color=CYAN, bold=True)
divider(slide, Inches(6.8), Inches(4.87), Inches(6.0), color=CYAN, thick=1)

env_items = [
    ('核函数', 'RBF (径向基函数)', BLUE),
    ('C参数', '1.0 (惩罚系数)', GREEN),
    ('γ参数', '0.5 (核宽度)', ORANGE),
    ('特征维度', '64维 (16bin × 4cell)', PURPLE),
    ('训练总量', '16,395张字符图片 (34+31类)', CYAN),
]
for i, (key, val, color) in enumerate(env_items):
    ey = Inches(5.05) + i * Inches(0.42)
    rbox(slide, Inches(6.8), ey, Inches(6.0), Inches(0.36), fill=CARD2, line_color=color, line_w=1)
    txt(slide, Inches(6.95), ey + Inches(0.04), Inches(1.8), Inches(0.26),
        key, size=11, color=color, bold=True)
    txt(slide, Inches(8.8), ey + Inches(0.04), Inches(3.8), Inches(0.26),
        val, size=10.5, color=GRAY)

# =========================================================
# SLIDE 9 — GUI & SYSTEM FEATURES (实际实现)
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title(slide, '04  系统功能与界面', 'Tkinter桌面端 · 双通道对比 · 多帧融合 · Web仪表盘')

# 左侧：GUI功能
rbox(slide, Inches(0.3), Inches(1.55), Inches(6.2), Inches(5.4), fill=CARD, line_color=BLUE, line_w=1.5)
txt(slide, Inches(0.5), Inches(1.7), Inches(5), Inches(0.4),
    'Tkinter桌面端功能', size=14, color=BLUE, bold=True)
divider(slide, Inches(0.5), Inches(2.12), Inches(5.8), color=BLUE, thick=1)

gui_features = [
    ('图片识别模式', '双线程并行执行形状+颜色通道定位', BLUE),
    ('摄像头识别模式', '2秒/帧采集，后台线程异步识别', CYAN),
    ('双通道对比展示', '左:原图 | 右上:形状定位 | 右下:颜色定位', GREEN),
    ('车牌颜色识别', '自动分类蓝牌/黄牌/绿牌并标记', ORANGE),
    ('多帧投票融合', '窗口=5帧, 阈值=3次, 提升摄像头稳定性', PURPLE),
    ('SQLite持久化', '每条识别结果存入vlpr_records.db', RED),
]
for i, (name, desc, color) in enumerate(gui_features):
    ny = Inches(2.3) + i * Inches(0.68)
    rbox(slide, Inches(0.5), ny, Inches(5.8), Inches(0.58), fill=CARD2, line_color=color, line_w=1)
    tag(slide, Inches(0.55), ny + Inches(0.13), name, color, fsize=9)
    txt(slide, Inches(1.7), ny + Inches(0.1), Inches(4.4), Inches(0.35),
        desc, size=10.5, color=GRAY)

# 右侧：Web仪表盘 + 数据推送
rbox(slide, Inches(6.8), Inches(1.55), Inches(6.2), Inches(2.5), fill=CARD, line_color=CYAN, line_w=1.5)
txt(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.4),
    'Web可视化仪表盘', size=14, color=CYAN, bold=True)
divider(slide, Inches(7.0), Inches(2.12), Inches(5.8), color=CYAN, thick=1)

dashboard_items = [
    '暗色主题 Chart.js 可视化图表',
    'WebSocket实时推送识别结果至仪表盘',
    '识别记录历史查询与统计分析',
    '车牌颜色分布、识别趋势等数据展示',
]
for i, item in enumerate(dashboard_items):
    circle(slide, Inches(7.0), Inches(2.3) + i * Inches(0.52), Inches(0.26), CYAN, str(i+1), fsize=9)
    txt(slide, Inches(7.38), Inches(2.3) + i * Inches(0.52), Inches(5.3), Inches(0.48),
        item, size=11, color=GRAY)

# Android部署
rbox(slide, Inches(6.8), Inches(4.3), Inches(6.2), Inches(2.65), fill=CARD, line_color=GREEN, line_w=1.5)
txt(slide, Inches(7.0), Inches(4.45), Inches(5), Inches(0.4),
    'Android移动端部署', size=14, color=GREEN, bold=True)
divider(slide, Inches(7.0), Inches(4.87), Inches(5.8), color=GREEN, thick=1)

android_items = [
    ('UI框架', 'Kivy (跨平台Python GUI框架)', BLUE),
    ('打包工具', 'buildozer 一键打包APK', GREEN),
    ('目标设备', 'Android 12 · 华为Mate 40 Pro · 骁龙865', ORANGE),
    ('核心功能', '相机拍照识别 · 结果实时显示', PURPLE),
]
for i, (key, val, color) in enumerate(android_items):
    ay = Inches(5.05) + i * Inches(0.48)
    rbox(slide, Inches(7.0), ay, Inches(5.8), Inches(0.4), fill=CARD2, line_color=color, line_w=1)
    txt(slide, Inches(7.15), ay + Inches(0.06), Inches(1.6), Inches(0.28),
        key, size=10.5, color=color, bold=True)
    txt(slide, Inches(8.8), ay + Inches(0.06), Inches(3.8), Inches(0.28),
        val, size=10.5, color=GRAY)

# =========================================================
# SLIDE 10 — TRAINING DATA & MODEL (实际训练数据)
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title(slide, '05  训练数据与模型', '16,395张字符图片 · 双SVM模型 · RBF核分类')

# 训练数据概览
rbox(slide, Inches(0.3), Inches(1.55), Inches(6.2), Inches(5.4), fill=CARD, line_color=BLUE, line_w=1.5)
txt(slide, Inches(0.5), Inches(1.7), Inches(5), Inches(0.4),
    '训练数据分布', size=14, color=BLUE, bold=True)
divider(slide, Inches(0.5), Inches(2.12), Inches(5.8), color=BLUE, thick=1)

# 中文省份
txt(slide, Inches(0.5), Inches(2.3), Inches(5.5), Inches(0.4),
    '中文省份训练集 (train/charsChinese/)', size=12, color=ORANGE, bold=True)

province_data = [
    '31个省份简称类别（川、鄂、赣、甘、贵、桂等）',
    '共3,232张汉字训练样本',
    '样本量分布：津(299) > 浙(257) > 沪(203) > 苏(219)',
    '部分稀有省份样本较少：藏(9)、新(19)',
    '归一化尺寸：20×20像素',
]
for i, pt in enumerate(province_data):
    py = Inches(2.72) + i * Inches(0.45)
    circle(slide, Inches(0.5), py + Pt(2), Inches(0.24), ORANGE, str(i+1), fsize=9)
    txt(slide, Inches(0.84), py, Inches(5.4), Inches(0.4),
        pt, size=10.5, color=GRAY)

# 字母数字
txt(slide, Inches(0.5), Inches(4.95), Inches(5.5), Inches(0.4),
    '字母数字训练集 (train/chars2/)', size=12, color=CYAN, bold=True)

alnum_data = [
    '34类：数字0-9 + 字母A-Z',
    '共13,163张训练样本',
    '数字"1"样本最多(1,154张)，字母"V"(80)、"Z"(98)较少',
    '字符倾斜校正(deskew)后提取HOG特征',
    '使用OpenCV cv2.ml.SVM_create()训练',
]
for i, pt in enumerate(alnum_data):
    py = Inches(5.37) + i * Inches(0.45)
    circle(slide, Inches(0.5), py + Pt(2), Inches(0.24), CYAN, str(i+1), fsize=9)
    txt(slide, Inches(0.84), py, Inches(5.4), Inches(0.4),
        pt, size=10.5, color=GRAY)

# 右侧：模型架构图
rbox(slide, Inches(6.8), Inches(1.55), Inches(6.2), Inches(5.4), fill=CARD, line_color=GREEN, line_w=1.5)
txt(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.4),
    'SVM识别流程', size=14, color=GREEN, bold=True)
divider(slide, Inches(7.0), Inches(2.12), Inches(5.8), color=GREEN, thick=1)

# 流程图
flow_steps = [
    ('车牌图像', DGRAY),
    ('投影分割\n7个字符', BLUE),
    ('补齐+resize\n20×20', CYAN),
    ('deskew\n倾斜校正', GREEN),
    ('HOG提取\n64维特征', ORANGE),
    ('SVM预测\nRBF核', PURPLE),
]
fw = Inches(1.5)
fy = Inches(2.35)
fh = Inches(0.85)
fsx0 = Inches(7.0)

for i, (label, color) in enumerate(flow_steps):
    fx = fsx0 + i * (fw + Inches(0.2))
    rbox(slide, fx, fy, fw, fh, fill=color)
    txt(slide, fx, fy, fw, fh, label, size=10, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER)
    if i < len(flow_steps)-1:
        arrow_right(slide, fx + fw, fy + Inches(0.3), Inches(0.18), Inches(0.22), DGRAY)

# 双模型分支
txt(slide, Inches(7.0), Inches(3.4), Inches(5.8), Inches(0.35),
    '首字符 → svmchinese.dat (31类省份)', size=11, color=ORANGE, bold=True)
txt(slide, Inches(7.0), Inches(3.75), Inches(5.8), Inches(0.35),
    '后6字符 → svm.dat (34类字母数字)', size=11, color=CYAN, bold=True)

# 关键技术参数
txt(slide, Inches(7.0), Inches(4.3), Inches(5.5), Inches(0.4),
    '关键技术参数', size=13, color=GREEN, bold=True)
divider(slide, Inches(7.0), Inches(4.7), Inches(5.8), color=GREEN, thick=1)

params = [
    ('字符尺寸', '20×20 像素'),
    ('HOG方向数', '16 bin'),
    ('Cell尺寸', '10×10 像素'),
    ('Cell数量', '4个 (2×2布局)'),
    ('特征维度', '16×4 = 64维'),
    ('核函数', 'RBF (高斯核)'),
    ('归一化', 'Hellinger + L2'),
    ('判尾过滤', '高宽比≥7的"1"剔除'),
]
for i, (k, v) in enumerate(params):
    col = i % 2
    row = i // 2
    px = Inches(7.0) + col * Inches(3.0)
    py = Inches(4.85) + row * Inches(0.45)
    txt(slide, px, py, Inches(1.4), Inches(0.35),
        k, size=10.5, color=DGRAY, bold=True)
    txt(slide, px + Inches(1.5), py, Inches(1.3), Inches(0.35),
        v, size=10.5, color=WHITE)

# =========================================================
# SLIDE 11 — CONCLUSION & OUTLOOK
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_title(slide, '06  总结与展望', '成果归纳 · 现存局限 · 优化方向')

# 左：成果
rbox(slide, Inches(0.3), Inches(1.55), Inches(6.0), Inches(5.4), fill=CARD, line_color=GREEN, line_w=1.5)
txt(slide, Inches(0.5), Inches(1.7), Inches(4), Inches(0.4),
    '研究成果', size=16, color=GREEN, bold=True)
divider(slide, Inches(0.5), Inches(2.12), Inches(5.6), color=GREEN, thick=1)

achievements = [
    ('双通道互补定位', '形状通道+颜色通道并行执行，结果互为验证，提升复杂场景鲁棒性'),
    ('HOG+SVM分类', '64维HOG特征+RBF-SVM双模型，中文31类+字母数字34类全覆盖'),
    ('多帧投票融合', '滑动窗口5帧+3次确认阈值，有效提升摄像头模式识别稳定性'),
    ('多平台部署', 'Tkinter桌面端 + Kivy Android端 + Chart.js Web仪表盘全覆盖'),
    ('完整数据管线', 'SQLite持久化 + WebSocket实时推送 + 历史查询统计分析'),
]
for i, (title, desc) in enumerate(achievements):
    ay = Inches(2.3) + i * Inches(0.92)
    circle(slide, Inches(0.55), ay + Pt(3), Inches(0.38), GREEN, str(i+1), fsize=14)
    txt(slide, Inches(1.1), ay, Inches(5.0), Inches(0.35),
        title, size=13, color=WHITE, bold=True)
    txt(slide, Inches(1.1), ay + Inches(0.38), Inches(5.0), Inches(0.5),
        desc, size=10.5, color=GRAY)

# 右：局限
rbox(slide, Inches(6.6), Inches(1.55), Inches(6.4), Inches(2.5), fill=CARD, line_color=RED, line_w=1.5)
txt(slide, Inches(6.8), Inches(1.7), Inches(4), Inches(0.4),
    '现存局限', size=16, color=RED, bold=True)
divider(slide, Inches(6.8), Inches(2.12), Inches(6.0), color=RED, thick=1)

limits = [
    '极端光照（强逆光/夜间）下车牌区域特征表达不足',
    '严重污损/遮挡车牌泛化能力有限',
    '部分稀有省份训练样本不足（如藏仅9张）',
    'SVM不输出概率值，无法评估识别置信度',
]
for i, lim in enumerate(limits):
    circle(slide, Inches(6.8), Inches(2.3) + i * Inches(0.55), Inches(0.28), RED, '!', fsize=10)
    txt(slide, Inches(7.2), Inches(2.3) + i * Inches(0.55), Inches(5.6), Inches(0.48),
        lim, size=11.5, color=GRAY)

# 优化方向
rbox(slide, Inches(6.6), Inches(4.3), Inches(6.4), Inches(2.65), fill=CARD, line_color=CYAN, line_w=1.5)
txt(slide, Inches(6.8), Inches(4.45), Inches(4), Inches(0.4),
    '未来优化方向', size=16, color=CYAN, bold=True)
divider(slide, Inches(6.8), Inches(4.87), Inches(6.0), color=CYAN, thick=1)

outlooks = [
    '引入深度学习(CNN/YOLO)提升复杂场景识别准确率',
    '数据增强扩充稀有类别样本，提升泛化能力',
    'Retinex理论光照校正，适应极端光照条件',
    '迁移学习适配新能源车牌等新类型',
    '轻量化模型部署至嵌入式设备(树莓派等)',
]
for i, out in enumerate(outlooks):
    circle(slide, Inches(6.8), Inches(5.05) + i * Inches(0.42), Inches(0.26), CYAN, '→', fsize=9)
    txt(slide, Inches(7.2), Inches(5.05) + i * Inches(0.42), Inches(5.6), Inches(0.38),
        out, size=11, color=GRAY)

# =========================================================
# SLIDE 12 — THANK YOU
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG)

box(slide, 0, 0, prs.slide_width, Inches(0.08), fill=BLUE)
box(slide, 0, Inches(0.08), prs.slide_width, Inches(0.04), fill=CYAN)

box(slide, Inches(4.2), Inches(1.0), Inches(5.0), Inches(5.0),
    fill=RGBColor(0x14, 0x2D, 0x6E), shape=MSO_SHAPE.OVAL).line.color.rgb = BLUE
box(slide, Inches(4.9), Inches(1.7), Inches(3.6), Inches(3.6),
    fill=RGBColor(0x1A, 0x38, 0x80), shape=MSO_SHAPE.OVAL).line.fill.background()

txt(slide, Inches(4.2), Inches(2.6), Inches(5.0), Inches(1.8),
    '谢谢聆听！\n敬请批评指正',
    size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(slide, Inches(0), Inches(6.0), prs.slide_width, Inches(0.7),
    '基于OpenCV的工程车车牌识别系统设计与实现  ·  李志杰  ·  2022b15080  ·  物联网工程  ·  2026届',
    size=13, color=GRAY, align=PP_ALIGN.CENTER)

divider(slide, Inches(2.5), Inches(5.88), Inches(8.3), color=BLUE, thick=2)

# 实际项目关键数据
nums = [
    ('16,395', '训练样本总量'),
    ('64维', 'HOG特征维度'),
    ('65类', '字符分类总数'),
    ('双通道', '互补定位架构'),
]
for i, (n, l) in enumerate(nums):
    nx = Inches(0.8) + i * Inches(3.1)
    rbox(slide, nx, Inches(6.75), Inches(2.6), Inches(0.65), fill=CARD, line_color=BLUE, line_w=1)
    txt(slide, nx, Inches(6.78), Inches(2.6), Inches(0.32),
        n, size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, nx, Inches(7.1), Inches(2.6), Inches(0.28),
        l, size=10, color=GRAY, align=PP_ALIGN.CENTER)

# =========================================================
# SAVE
# =========================================================
out = r'C:\Users\HUAWEI\Desktop\FinalDesign\Python_VLPR-master\file\答辩PPT_李志杰_2022b15080_v3.pptx'
prs.save(out)
print(f'Saved: {out}')
print(f'Total slides: {len(prs.slides)}')
