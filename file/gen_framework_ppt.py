"""
论文整体框架图生成脚本 — 基于实际项目产物
参考用户提供的框架图样式，绘制六章结构的论文整体框架
所有内容基于项目实际代码实现：OpenCV + HOG + SVM + Tkinter
"""
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# =========================================================
# 1. PIL绘制高清框架图
# =========================================================
W, H = 1920, 1080
NAVY = (13, 43, 94)
WHITE = (255, 255, 255)
LIGHT_BG = (245, 247, 250)
DARK_TEXT = (26, 26, 46)
MID_TEXT = (74, 85, 104)
CARD_BORDER = (13, 43, 94)
SEPARATOR = (180, 190, 210)

img = Image.new('RGB', (W, H), LIGHT_BG)
draw = ImageDraw.Draw(img)

# 加载字体
def load_font(path, size):
    try:
        return ImageFont.truetype(path, size)
    except:
        return ImageFont.load_default()

font_header = load_font("C:/Windows/Fonts/msyhbd.ttc", 20)
font_card_title = load_font("C:/Windows/Fonts/msyhbd.ttc", 15)
font_card_desc = load_font("C:/Windows/Fonts/msyh.ttc", 13)
font_card_small = load_font("C:/Windows/Fonts/msyh.ttc", 11)
font_bottom = load_font("C:/Windows/Fonts/msyhbd.ttc", 19)
font_bottom_sub = load_font("C:/Windows/Fonts/msyh.ttc", 13)

# 圆角矩形绘制
def draw_rounded_rect(draw, xy, radius, fill, outline=None, width=2):
    x1, y1, x2, y2 = xy
    # 主体矩形
    draw.rectangle([x1+radius, y1, x2-radius, y2], fill=fill)
    draw.rectangle([x1, y1+radius, x2, y2-radius], fill=fill)
    # 四个圆角（使用椭圆扇形）
    draw.pieslice([x1, y1, x1+2*radius, y1+2*radius], 180, 270, fill=fill)
    draw.pieslice([x2-2*radius, y1, x2, y1+2*radius], 270, 360, fill=fill)
    draw.pieslice([x1, y2-2*radius, x1+2*radius, y2], 90, 180, fill=fill)
    draw.pieslice([x2-2*radius, y2-2*radius, x2, y2], 0, 90, fill=fill)
    # 描边
    if outline:
        draw.arc([x1, y1, x1+2*radius, y1+2*radius], 180, 270, fill=outline, width=width)
        draw.arc([x2-2*radius, y1, x2, y1+2*radius], 270, 360, fill=outline, width=width)
        draw.arc([x1, y2-2*radius, x1+2*radius, y2], 90, 180, fill=outline, width=width)
        draw.arc([x2-2*radius, y2-2*radius, x2, y2], 0, 90, fill=outline, width=width)
        draw.line([(x1+radius, y1), (x2-radius, y1)], fill=outline, width=width)
        draw.line([(x1+radius, y2), (x2-radius, y2)], fill=outline, width=width)
        draw.line([(x1, y1+radius), (x1, y2-radius)], fill=outline, width=width)
        draw.line([(x2, y1+radius), (x2, y2-radius)], fill=outline, width=width)

# 绘制章节标题
def draw_chapter_header(text, y):
    bbox = draw.textbbox((0, 0), text, font=font_header)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    pad_x, pad_y = 36, 10
    hw = tw + 2 * pad_x
    hh = th + 2 * pad_y + 2
    hx = (W - hw) // 2
    draw_rounded_rect(draw, (hx, y, hx + hw, y + hh), 8, NAVY)
    draw.text((hx + pad_x, y + pad_y), text, fill=WHITE, font=font_header)
    return hh

# 绘制内容卡片
def draw_cards(items, y, card_h=68):
    n = len(items)
    margin = 75
    gap = 20 if n <= 3 else (16 if n == 4 else 14)
    available_w = W - 2 * margin
    card_w = (available_w - (n - 1) * gap) // n
    start_x = margin
    
    for i, (title, desc) in enumerate(items):
        cx = start_x + i * (card_w + gap)
        # 卡片背景+边框
        draw_rounded_rect(draw, (cx, y, cx + card_w, y + card_h), 10, WHITE, CARD_BORDER, 2)
        
        # 标题
        draw.text((cx + 14, y + 10), title, fill=DARK_TEXT, font=font_card_title)
        
        # 描述文字（自动换行）
        if desc:
            words = desc
            max_width = card_w - 24
            lines = []
            current_line = ""
            for char in words:
                test = current_line + char
                bb = draw.textbbox((0, 0), test, font=font_card_desc)
                if bb[2] - bb[0] > max_width and current_line:
                    lines.append(current_line)
                    current_line = char
                else:
                    current_line = test
            if current_line:
                lines.append(current_line)
            
            # 最多显示2行
            for li, line in enumerate(lines[:2]):
                draw.text((cx + 14, y + 32 + li * 18), line, fill=MID_TEXT, font=font_card_desc)

# 绘制虚线分隔符
def draw_dashed_line(y):
    for x in range(70, W - 70, 16):
        draw.line([(x, y), (x + 9, y)], fill=SEPARATOR, width=1)

# 绘制向下箭头
def draw_down_arrow(y):
    ax = W // 2
    ay = y
    # 箭头主体（小矩形）
    draw.polygon([(ax - 7, ay), (ax + 7, ay), (ax, ay + 14)], fill=NAVY)

# =========================================================
# 六章数据结构 — 基于实际项目产物
# =========================================================
sections = [
    ("第1章 绪论", [
        ("研究背景与意义", "工程车复杂场景·人工登记低效·智慧工地需求"),
        ("国内外研究现状", "传统图像处理·OpenCV方法·车牌识别研究现状"),
        ("研究内容与技术路线", "图像采集·双通道定位·HOG+SVM识别·多平台部署"),
    ]),
    ("第2章 系统相关技术基础", [
        ("OpenCV图像处理", "高斯滤波·灰度化·形态学操作·边缘检测·轮廓筛选"),
        ("HOG特征提取", "Sobel梯度·16方向bin·4个cell·64维特征向量"),
        ("SVM支持向量机", "RBF核函数·C-SVC分类·中文31类·字母数字34类"),
        ("车牌定位算法", "双通道互补定位·形状通道·颜色通道·HSV空间筛选"),
        ("字符识别算法", "HOG特征·RBF-SVM预测·双模型并行·投影法分割"),
    ]),
    ("第3章 系统需求分析与总体设计", [
        ("功能需求分析", "图片选择·摄像头采集·车牌定位·字符识别·结果展示"),
        ("性能需求", "双通道验证鲁棒性·推理快速·模型轻量·多平台兼容"),
        ("系统总体架构", "图像输入·双通道定位·SVM识别·多平台结果输出"),
        ("模块划分", "图像预处理·双通道车牌定位·字符分割识别·GUI与数据展示"),
    ]),
    ("第4章 系统详细设计与实现", [
        ("图像预处理模块", "高斯模糊·灰度化·形态学开闭·Otsu二值化·Canny边缘"),
        ("双通道车牌定位", "形状通道:轮廓筛选·颜色通道:HSV掩码·仿射变换校正"),
        ("字符识别模块", "HOG 64维特征·RBF-SVM双模型·31省+34类·16,395张训练"),
        ("GUI可视化模块", "Tkinter界面·双通道对比·多帧融合·SQLite记录·WebSocket推送"),
    ]),
    ("第5章 系统测试与结果分析", [
        ("测试环境与方案", "OpenCV+Python·16,395张训练样本·25张测试图片"),
        ("分模块测试", "预处理·形状定位·颜色定位·HOG特征·SVM预测·GUI交互"),
        ("系统整体测试", "本地图片模式·摄像头模式·多帧投票融合·Web仪表盘"),
        ("问题分析与优化", "极端光照·稀有省份样本不足·未来可引入深度学习"),
    ]),
    ("第6章 总结与展望", [
        ("研究成果总结", "双通道互补定位·HOG+SVM识别·多平台部署·可视化仪表盘"),
        ("不足与未来展望", "深度学习提升准确率·数据增强扩充样本·嵌入式部署优化"),
    ]),
]

# =========================================================
# 绘制布局
# =========================================================
current_y = 22

for idx, (header, cards) in enumerate(sections):
    # 章节标题
    hh = draw_chapter_header(header, current_y)
    current_y += hh + 10
    
    # 内容卡片
    draw_cards(cards, current_y, card_h=72 if len(cards) <= 3 else 66)
    current_y += 76 if len(cards) <= 3 else 70
    
    # 分隔线 + 箭头（最后一张后面不加）
    if idx < len(sections) - 1:
        draw_dashed_line(current_y + 6)
        current_y += 18
        draw_down_arrow(current_y)
        current_y += 22

# =========================================================
# 底部总结栏
# =========================================================
bottom_y = current_y + 10
bar_h = 52
# 深蓝色背景条
draw_rounded_rect(draw, (70, bottom_y, W - 70, bottom_y + bar_h), 10, NAVY)

# 左侧图标（圆形箭头）
cx, cy = 110, bottom_y + bar_h // 2
r = 16
draw.ellipse([cx - r, cy - r, cx + r, cy + r], outline=WHITE, width=2)
# 箭头符号
draw.text((cx - 7, cy - 9), "⟳", fill=WHITE, font=font_bottom)

# 底部文字
draw.text((145, bottom_y + 12), "基于OpenCV的工程车车牌识别系统设计与实现", fill=WHITE, font=font_bottom)
draw.text((145, bottom_y + 32), "OpenCV图像处理 + HOG特征提取 + SVM分类识别 + 多平台部署", fill=(180, 195, 220), font=font_bottom_sub)

# =========================================================
# 保存PNG
# =========================================================
png_path = r"C:\Users\HUAWEI\Desktop\FinalDesign\Python_VLPR-master\file\论文整体框架图_v3.png"
img.save(png_path, quality=95)
print(f"PNG saved: {png_path}")

# =========================================================
# 2. 嵌入PPT
# =========================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# 白色背景
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# 添加图片，铺满整个幻灯片
left = Inches(0)
top = Inches(0)
pic = slide.shapes.add_picture(png_path, left, top, width=Inches(13.333), height=Inches(7.5))

# 保存PPT
ppt_path = r"C:\Users\HUAWEI\Desktop\FinalDesign\Python_VLPR-master\file\论文整体框架_v3.pptx"
prs.save(ppt_path)
print(f"PPT saved: {ppt_path}")
