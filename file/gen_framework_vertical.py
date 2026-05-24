"""
论文整体框架图生成脚本 — 竖屏版本（参考图规格）
基于实际项目产物：OpenCV + HOG + SVM + Tkinter
画布规格：1080×1920 竖屏，深蓝主题
"""
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os, math

# =========================================================
# 画布与配色（参考图风格）
# =========================================================
W, H = 1080, 1920
MARGIN = 50

NAVY = (27, 58, 120)           # 章节标题栏深蓝
CARD_BORDER = (52, 103, 168)   # 卡片边框蓝
WHITE = (255, 255, 255)
LIGHT_BG = (240, 244, 248)     # 浅蓝灰背景
DARK_TEXT = (30, 40, 60)       # 卡片标题深色
MID_TEXT = (80, 95, 115)       # 描述文字灰色
LIGHT_ICON_BG = (220, 230, 245) # 图标背景浅蓝
ARROW_COLOR = (27, 58, 120)
BOTTOM_BAR = (27, 58, 120)
DASH_COLOR = (160, 175, 195)

img = Image.new('RGB', (W, H), LIGHT_BG)
draw = ImageDraw.Draw(img)

# 字体加载
def load_font(path, size):
    try:
        return ImageFont.truetype(path, size)
    except:
        return ImageFont.load_default()

font_header = load_font("C:/Windows/Fonts/msyhbd.ttc", 22)
font_card_title = load_font("C:/Windows/Fonts/msyhbd.ttc", 14)
font_card_desc = load_font("C:/Windows/Fonts/msyh.ttc", 12)
font_bottom = load_font("C:/Windows/Fonts/msyhbd.ttc", 18)
font_bottom_sub = load_font("C:/Windows/Fonts/msyh.ttc", 13)

# =========================================================
# 绘制辅助函数
# =========================================================
def draw_rounded_rect(draw, xy, radius, fill, outline=None, width=2):
    x1, y1, x2, y2 = xy
    draw.rectangle([x1+radius, y1, x2-radius, y2], fill=fill)
    draw.rectangle([x1, y1+radius, x2, y2-radius], fill=fill)
    draw.pieslice([x1, y1, x1+2*radius, y1+2*radius], 180, 270, fill=fill)
    draw.pieslice([x2-2*radius, y1, x2, y1+2*radius], 270, 360, fill=fill)
    draw.pieslice([x1, y2-2*radius, x1+2*radius, y2], 90, 180, fill=fill)
    draw.pieslice([x2-2*radius, y2-2*radius, x2, y2], 0, 90, fill=fill)
    if outline:
        draw.arc([x1, y1, x1+2*radius, y1+2*radius], 180, 270, fill=outline, width=width)
        draw.arc([x2-2*radius, y1, x2, y1+2*radius], 270, 360, fill=outline, width=width)
        draw.arc([x1, y2-2*radius, x1+2*radius, y2], 90, 180, fill=outline, width=width)
        draw.arc([x2-2*radius, y2-2*radius, x2, y2], 0, 90, fill=outline, width=width)
        draw.line([(x1+radius, y1), (x2-radius, y1)], fill=outline, width=width)
        draw.line([(x1+radius, y2), (x2-radius, y2)], fill=outline, width=width)
        draw.line([(x1, y1+radius), (x1, y2-radius)], fill=outline, width=width)
        draw.line([(x2, y1+radius), (x2, y2-radius)], fill=outline, width=width)

# 图标绘制函数
def draw_icon(draw, icon_type, cx, cy, size=14):
    """在指定中心绘制简单几何图标"""
    s = size
    if icon_type == "doc":  # 文档/列表
        draw.rectangle([cx-s, cy-s, cx+s, cy+s], outline=CARD_BORDER, width=2)
        for i in range(-s//2, s//2+1, 5):
            draw.line([(cx-s+3, cy+i), (cx+s-3, cy+i)], fill=CARD_BORDER, width=1)
    elif icon_type == "globe":  # 地球/网络
        draw.ellipse([cx-s, cy-s, cx+s, cy+s], outline=CARD_BORDER, width=2)
        draw.line([(cx, cy-s), (cx, cy+s)], fill=CARD_BORDER, width=1)
        draw.line([(cx-s, cy), (cx+s, cy)], fill=CARD_BORDER, width=1)
    elif icon_type == "target":  # 靶心/目标
        draw.ellipse([cx-s, cy-s, cx+s, cy+s], outline=CARD_BORDER, width=2)
        draw.ellipse([cx-s//2, cy-s//2, cx+s//2, cy+s//2], outline=CARD_BORDER, width=1)
        draw.line([(cx-2, cy), (cx+2, cy)], fill=CARD_BORDER, width=2)
        draw.line([(cx, cy-2), (cx, cy+2)], fill=CARD_BORDER, width=2)
    elif icon_type == "image":  # 图像/照片
        draw.rectangle([cx-s, cy-s, cx+s, cy+s], outline=CARD_BORDER, width=2)
        draw.polygon([(cx-s+2, cy+s-2), (cx-2, cy+2), (cx+4, cy+s-6)], fill=CARD_BORDER)
        draw.ellipse([cx+2, cy-s+3, cx+6, cy-s+7], fill=CARD_BORDER)
    elif icon_type == "grid":  # 网格/HOG
        for i in range(-s, s+1, s//2):
            draw.line([(cx-s, cy+i), (cx+s, cy+i)], fill=CARD_BORDER, width=1)
            draw.line([(cx+i, cy-s), (cx+i, cy+s)], fill=CARD_BORDER, width=1)
    elif icon_type == "brain":  # SVM/ML
        draw.ellipse([cx-s, cy-s//2, cx+s, cy+s//2], outline=CARD_BORDER, width=2)
        draw.line([(cx-s//2, cy), (cx, cy-s//3), (cx+s//2, cy)], fill=CARD_BORDER, width=1)
        draw.line([(cx-s//2, cy), (cx, cy+s//3), (cx+s//2, cy)], fill=CARD_BORDER, width=1)
    elif icon_type == "car":  # 车牌/定位
        draw.rectangle([cx-s, cy-s//3, cx+s, cy+s//3], outline=CARD_BORDER, width=2, fill=LIGHT_ICON_BG)
        draw.ellipse([cx-s+3, cy+s//3, cx-s+8, cy+s//3+5], outline=CARD_BORDER, width=1)
        draw.ellipse([cx+s-8, cy+s//3, cx+s-3, cy+s//3+5], outline=CARD_BORDER, width=1)
    elif icon_type == "text":  # 文字/识别
        draw.rectangle([cx-s, cy-s, cx+s, cy+s], outline=CARD_BORDER, width=2)
        draw.text((cx-5, cy-6), "A", fill=CARD_BORDER, font=font_card_title)
    elif icon_type == "list":  # 清单/需求
        for i, y_off in enumerate([-6, 0, 6]):
            draw.rectangle([cx-s+2, cy+y_off-1, cx-s+6, cy+y_off+1], fill=CARD_BORDER)
            draw.line([(cx-s+8, cy+y_off), (cx+s-2, cy+y_off)], fill=CARD_BORDER, width=1)
    elif icon_type == "gauge":  # 性能/仪表盘
        draw.arc([cx-s, cy-s, cx+s, cy+s], 0, 180, fill=CARD_BORDER, width=2)
        draw.line([(cx, cy), (cx+6, cy-8)], fill=CARD_BORDER, width=2)
    elif icon_type == "layers":  # 架构/分层
        for i, y_off in enumerate([-8, 0, 8]):
            w = s - abs(i-1)*3
            draw.rectangle([cx-w, cy+y_off-2, cx+w, cy+y_off+2], outline=CARD_BORDER, width=1)
    elif icon_type == "puzzle":  # 模块/拼图
        draw.rectangle([cx-s+2, cy-s+2, cx+s-2, cy+s-2], outline=CARD_BORDER, width=2)
        draw.rectangle([cx-2, cy-s, cx+2, cy-s+4], outline=CARD_BORDER, width=1, fill=LIGHT_BG)
        draw.rectangle([cx-2, cy+s-4, cx+2, cy+s], outline=CARD_BORDER, width=1, fill=LIGHT_BG)
    elif icon_type == "filter":  # 滤波/预处理
        draw.polygon([(cx-s, cy-s), (cx+s, cy-s), (cx, cy+s)], outline=CARD_BORDER, width=2)
        draw.line([(cx-s//2, cy), (cx+s//2, cy)], fill=CARD_BORDER, width=1)
    elif icon_type == "split":  # 双通道
        draw.rectangle([cx-s, cy-s, cx, cy+s], outline=CARD_BORDER, width=2)
        draw.rectangle([cx, cy-s, cx+s, cy+s], outline=CARD_BORDER, width=2)
        draw.line([(cx, cy-s), (cx, cy+s)], fill=CARD_BORDER, width=1)
    elif icon_type == "screen":  # GUI/屏幕
        draw.rectangle([cx-s, cy-s, cx+s, cy+s-4], outline=CARD_BORDER, width=2)
        draw.rectangle([cx-3, cy+s-4, cx+3, cy+s], fill=CARD_BORDER)
    elif icon_type == "gear":  # 环境/设置
        n = 8
        r_out, r_in = s, s//2
        pts_out = []
        pts_in = []
        for i in range(n*2):
            angle = math.pi * 2 * i / (n*2) - math.pi/2
            r = r_out if i % 2 == 0 else r_in
            pts_out.append((cx + r*math.cos(angle), cy + r*math.sin(angle)))
        draw.polygon(pts_out, outline=CARD_BORDER, width=2)
        draw.ellipse([cx-3, cy-3, cx+3, cy+3], fill=CARD_BORDER)
    elif icon_type == "test":  # 测试
        draw.rectangle([cx-s//2-2, cy-s, cx+s//2+2, cy+s], outline=CARD_BORDER, width=2)
        draw.line([(cx, cy-s), (cx, cy-s+5)], fill=CARD_BORDER, width=2)
        draw.line([(cx-4, cy-s+8), (cx+4, cy-s+8)], fill=CARD_BORDER, width=1)
    elif icon_type == "chart":  # 图表/分析
        draw.line([(cx-s, cy+s), (cx-s, cy-s)], fill=CARD_BORDER, width=2)
        draw.line([(cx-s, cy+s), (cx+s, cy+s)], fill=CARD_BORDER, width=2)
        for x, h in [(-s+4, -s+8), (0, -s+4), (s-4, -s+12)]:
            draw.rectangle([cx+x-3, cy+h, cx+x+3, cy+s-2], fill=CARD_BORDER)
    elif icon_type == "system":  # 系统/整体
        draw.rectangle([cx-s, cy-s, cx+s, cy+s], outline=CARD_BORDER, width=2)
        draw.ellipse([cx-s//2, cy-s//2, cx+s//2, cy+s//2], outline=CARD_BORDER, width=1)
        for pt in [(cx, cy-s), (cx, cy+s), (cx-s, cy), (cx+s, cy)]:
            draw.ellipse([pt[0]-2, pt[1]-2, pt[0]+2, pt[1]+2], fill=CARD_BORDER)
    elif icon_type == "wrench":  # 优化/扳手
        draw.line([(cx-8, cy-8), (cx+8, cy+8)], fill=CARD_BORDER, width=3)
        draw.ellipse([cx+4, cy+4, cx+10, cy+10], outline=CARD_BORDER, width=2)
        draw.ellipse([cx-10, cy-10, cx-4, cy-4], outline=CARD_BORDER, width=2)
    elif icon_type == "star":  # 成果/星星
        pts = []
        for i in range(10):
            angle = math.pi * 2 * i / 10 - math.pi / 2
            r = s if i % 2 == 0 else s // 2
            pts.append((cx + r * math.cos(angle), cy + r * math.sin(angle)))
        draw.polygon(pts, outline=CARD_BORDER, width=2)
    elif icon_type == "scope":  # 展望/望远镜
        draw.rectangle([cx-s, cy-3, cx+s, cy+3], outline=CARD_BORDER, width=2, fill=CARD_BORDER)
        draw.rectangle([cx+s-2, cy-8, cx+s+4, cy+8], outline=CARD_BORDER, width=2)
        draw.line([(cx-s-4, cy+3), (cx-s, cy+3)], fill=CARD_BORDER, width=2)

# =========================================================
# 绘制章节标题栏
# =========================================================
def draw_chapter_header(text, y):
    bbox = draw.textbbox((0, 0), text, font=font_header)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    pad_x, pad_y = 40, 10
    hw = W - 2 * MARGIN  # 标题栏宽度 = 画布宽 - 两边边距
    hh = th + 2 * pad_y + 2
    hx = MARGIN
    draw_rounded_rect(draw, (hx, y, hx + hw, y + hh), 8, NAVY)
    # 文字居中
    tx = hx + (hw - tw) // 2
    draw.text((tx, y + pad_y), text, fill=WHITE, font=font_header)
    return hh

# =========================================================
# 绘制内容卡片（图标+标题）
# =========================================================
def draw_cards(items, y, card_h=80):
    n = len(items)
    gap = 16 if n <= 3 else (14 if n == 4 else 12)
    available_w = W - 2 * MARGIN
    card_w = (available_w - (n - 1) * gap) // n
    start_x = MARGIN
    
    for i, (title, desc, icon_type) in enumerate(items):
        cx = start_x + i * (card_w + gap)
        cy = y
        
        # 卡片背景+边框
        draw_rounded_rect(draw, (cx, cy, cx + card_w, cy + card_h), 10, WHITE, CARD_BORDER, 2)
        
        # 图标（居中在顶部区域）
        icon_cx = cx + card_w // 2
        icon_cy = cy + 24
        # 绘制浅蓝色圆形背景
        r_bg = 18
        draw.ellipse([icon_cx-r_bg, icon_cy-r_bg, icon_cx+r_bg, icon_cy+r_bg], fill=LIGHT_ICON_BG)
        draw_icon(draw, icon_type, icon_cx, icon_cy, size=11)
        
        # 标题（图标下方，居中）
        tb = draw.textbbox((0, 0), title, font=font_card_title)
        tw = tb[2] - tb[0]
        tx = cx + (card_w - tw) // 2
        draw.text((tx, cy + 48), title, fill=DARK_TEXT, font=font_card_title)

# =========================================================
# 绘制虚线分隔 + 向下箭头
# =========================================================
def draw_dashed_line(y):
    for x in range(MARGIN, W - MARGIN, 14):
        draw.line([(x, y), (x + 7, y)], fill=DASH_COLOR, width=1)

def draw_down_arrow(y):
    ax = W // 2
    ay = y
    # 箭头主体（填充三角形）
    draw.polygon([(ax - 8, ay), (ax + 8, ay), (ax, ay + 16)], fill=ARROW_COLOR)

# =========================================================
# 六章数据结构 — 基于实际项目产物（带图标类型）
# =========================================================
sections = [
    ("第1章 绪论", [
        ("研究背景与意义", "工程车场景复杂·人工登记低效\n智慧工地数字化管理需求", "doc"),
        ("国内外研究现状", "传统图像处理方法\nOpenCV与车牌识别研究现状", "globe"),
        ("研究内容与技术路线", "图像采集·双通道定位\nHOG+SVM识别·多平台部署", "target"),
    ]),
    ("第2章 系统相关技术基础", [
        ("OpenCV图像处理", "高斯滤波·灰度化·形态学\n边缘检测·轮廓筛选", "image"),
        ("HOG特征提取", "Sobel梯度·16方向bin\n4个cell·64维特征向量", "grid"),
        ("SVM支持向量机", "RBF核函数·C-SVC分类\n中文31类·字母数字34类", "brain"),
        ("车牌定位算法", "双通道互补定位·形状通道\n颜色通道·HSV空间筛选", "car"),
        ("字符识别算法", "HOG特征·RBF-SVM预测\n双模型并行·投影法分割", "text"),
    ]),
    ("第3章 系统需求分析与总体设计", [
        ("功能需求分析", "图片选择·摄像头采集\n车牌定位·字符识别", "list"),
        ("性能需求", "双通道验证鲁棒性\n推理快速·模型轻量", "gauge"),
        ("系统总体架构", "图像输入·双通道定位\nSVM识别·多平台输出", "layers"),
        ("模块划分", "图像预处理·车牌定位\n字符识别·GUI与数据展示", "puzzle"),
    ]),
    ("第4章 系统详细设计与实现", [
        ("图像预处理模块", "高斯模糊·灰度化·形态学\nOtsu二值化·Canny边缘", "filter"),
        ("双通道车牌定位", "形状通道:轮廓筛选\n颜色通道:HSV掩码·仿射变换", "split"),
        ("字符识别模块", "HOG 64维特征·RBF-SVM\n31省+34类·16,395张训练", "text"),
        ("GUI可视化模块", "Tkinter界面·双通道对比\n多帧融合·SQLite记录", "screen"),
    ]),
    ("第5章 系统测试与结果分析", [
        ("测试环境与方案", "OpenCV+Python环境\n16,395张训练·25张测试", "gear"),
        ("分模块测试", "预处理·形状定位·颜色定位\nHOG特征·SVM预测", "test"),
        ("系统整体测试", "本地图片·摄像头模式\n多帧投票融合·Web仪表盘", "system"),
        ("问题分析与优化", "极端光照·稀有省份样本\n未来引入深度学习", "wrench"),
    ]),
    ("第6章 总结与展望", [
        ("研究成果总结", "双通道互补定位·HOG+SVM\n多平台部署·可视化仪表盘", "star"),
        ("不足与未来展望", "深度学习提升准确率\n数据增强·嵌入式部署优化", "scope"),
    ]),
]

# =========================================================
# 绘制布局
# =========================================================
current_y = 30

for idx, (header, cards) in enumerate(sections):
    # 章节标题
    hh = draw_chapter_header(header, current_y)
    current_y += hh + 14
    
    # 内容卡片（根据卡片数量调整高度）
    card_h = 82 if len(cards) <= 3 else 76
    draw_cards(cards, current_y, card_h=card_h)
    current_y += card_h + 10
    
    # 分隔线 + 箭头（最后一张后面不加）
    if idx < len(sections) - 1:
        draw_dashed_line(current_y + 6)
        current_y += 18
        draw_down_arrow(current_y)
        current_y += 26

# =========================================================
# 底部总结栏（匹配参考图风格）
# =========================================================
bottom_y = current_y + 14
bar_h = 56

# 深蓝色背景条（圆角）
draw_rounded_rect(draw, (MARGIN, bottom_y, W - MARGIN, bottom_y + bar_h), 10, BOTTOM_BAR)

# 左侧循环箭头图标
cx_icon, cy_icon = MARGIN + 38, bottom_y + bar_h // 2
r = 16
draw.ellipse([cx_icon-r, cy_icon-r, cx_icon+r, cy_icon+r], outline=WHITE, width=2)
# 绘制旋转箭头
for angle in range(0, 360, 120):
    rad = math.radians(angle - 90)
    x1 = cx_icon + 8 * math.cos(rad)
    y1 = cy_icon + 8 * math.sin(rad)
    x2 = cx_icon + 12 * math.cos(rad)
    y2 = cy_icon + 12 * math.sin(rad)
    draw.line([(x1, y1), (x2, y2)], fill=WHITE, width=2)

# 底部主标题
title_text = "基于OpenCV的工程车车牌识别系统设计与实现"
tb = draw.textbbox((0, 0), title_text, font=font_bottom)
tw = tb[2] - tb[0]
tx = MARGIN + 68
draw.text((tx, bottom_y + 10), title_text, fill=WHITE, font=font_bottom)

# 底部副标题
sub_text = "OpenCV图像处理 + HOG特征提取 + SVM分类识别 + 多平台部署"
draw.text((tx, bottom_y + 32), sub_text, fill=(180, 195, 220), font=font_bottom_sub)

# =========================================================
# 保存PNG
# =========================================================
png_path = r"C:\Users\HUAWEI\Desktop\FinalDesign\Python_VLPR-master\file\论文整体框架图_竖屏_v3.png"
img.save(png_path, quality=95)
print(f"PNG saved: {png_path}")

# =========================================================
# 嵌入PPT（竖屏比例 9:16）
# =========================================================
prs = Presentation()
prs.slide_width = Inches(7.5)
prs.slide_height = Inches(13.333)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# 白色背景
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# 添加图片，铺满整个幻灯片
left = Inches(0)
top = Inches(0)
pic = slide.shapes.add_picture(png_path, left, top, width=Inches(7.5), height=Inches(13.333))

# 保存PPT
ppt_path = r"C:\Users\HUAWEI\Desktop\FinalDesign\Python_VLPR-master\file\论文整体框架_竖屏_v3.pptx"
prs.save(ppt_path)
print(f"PPT saved: {ppt_path}")
