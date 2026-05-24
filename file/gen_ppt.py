from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

# === Color Palette ===
BG_DARK = RGBColor(0x1B, 0x1F, 0x3B)
BG_CARD = RGBColor(0x24, 0x29, 0x4A)
ACCENT_BLUE = RGBColor(0x4E, 0x9A, 0xFD)
ACCENT_CYAN = RGBColor(0x36, 0xD7, 0xB7)
ACCENT_ORANGE = RGBColor(0xFD, 0xA0, 0x5E)
ACCENT_PURPLE = RGBColor(0x9B, 0x6D, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xD1)
DIM_GRAY = RGBColor(0x6B, 0x72, 0x94)

# === Background ===
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = BG_DARK

# === Helper Functions ===
def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_arrow_right(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_circle_icon(slide, left, top, size, fill_color, text, font_size=16):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    return shape

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=10, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_before = Pt(0)
        p.space_after = Pt(0)
    return txBox

# =============================================
# TITLE AREA
# =============================================
add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.65),
            '技术路线', font_size=32, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.4),
            '\u201c数据驱动 + 模型优化 + 系统集成\u201d 三阶段技术路线',
            font_size=15, color=LIGHT_GRAY)

line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(3.5), Pt(3))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = ACCENT_BLUE
line_shape.line.fill.background()

# =============================================
# PHASE LABELS (Top Row)
# =============================================
phases = [
    ('01', '数据驱动', ACCENT_BLUE, 'CCPD数据集 -> 精细化预处理'),
    ('02', '模型优化', ACCENT_CYAN, 'OpenCV定位 + 字符神经网络'),
    ('03', '系统集成', ACCENT_ORANGE, '跨平台APP + 嵌入式部署'),
]

phase_start_x = Inches(0.8)
phase_gap = Inches(4.1)
phase_width = Inches(3.8)

for i, (num, name, color, subtitle) in enumerate(phases):
    x = phase_start_x + i * phase_gap
    y = Inches(1.65)
    add_circle_icon(slide, x, y, Inches(0.48), color, num, font_size=16)
    add_textbox(slide, x + Inches(0.6), y - Pt(2), Inches(2), Inches(0.48),
                name, font_size=22, color=color, bold=True)
    add_textbox(slide, x + Inches(0.6), y + Inches(0.35), Inches(3), Inches(0.35),
                subtitle, font_size=11, color=DIM_GRAY)

for i in range(2):
    x = phase_start_x + phase_width + i * phase_gap + Inches(0.15)
    add_arrow_right(slide, x, Inches(1.85), Inches(0.22), Inches(0.18), DIM_GRAY)

# =============================================
# PHASE 1: DATA DRIVEN
# =============================================
card_y = Inches(2.5)
card_h = Inches(2.1)

add_rounded_rect(slide, phase_start_x, card_y, phase_width, card_h, BG_CARD, ACCENT_BLUE)

items_phase1 = [
    ('CCPD数据集', '裁剪、灰度化、尺寸归一化'),
    ('数据增强', '旋转、噪声、亮度调节'),
    ('样本构建', '适用于复杂场景训练样本'),
]

for j, (title, desc) in enumerate(items_phase1):
    item_y = card_y + Inches(0.18) + j * Inches(0.65)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, phase_start_x + Inches(0.2), item_y + Pt(6), Pt(8), Pt(8))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_BLUE
    dot.line.fill.background()
    add_textbox(slide, phase_start_x + Inches(0.45), item_y - Pt(2), Inches(3), Inches(0.28),
                title, font_size=13, color=WHITE, bold=True)
    add_textbox(slide, phase_start_x + Inches(0.45), item_y + Inches(0.24), Inches(3.2), Inches(0.35),
                desc, font_size=11, color=LIGHT_GRAY)

# =============================================
# PHASE 2: MODEL OPTIMIZATION
# =============================================
x2 = phase_start_x + phase_gap
add_rounded_rect(slide, x2, card_y, phase_width, card_h, BG_CARD, ACCENT_CYAN)

items_phase2 = [
    ('OpenCV定位', '级联分类器实现车牌区域初定位'),
    ('字符识别网络', '无分割方案直接输出完整字符串'),
    ('模型对比调优', 'YOLOv3/SSD对比，平衡准确率与实时性'),
]

for j, (title, desc) in enumerate(items_phase2):
    item_y = card_y + Inches(0.18) + j * Inches(0.65)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x2 + Inches(0.2), item_y + Pt(6), Pt(8), Pt(8))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_CYAN
    dot.line.fill.background()
    add_textbox(slide, x2 + Inches(0.45), item_y - Pt(2), Inches(3), Inches(0.28),
                title, font_size=13, color=WHITE, bold=True)
    add_textbox(slide, x2 + Inches(0.45), item_y + Inches(0.24), Inches(3.2), Inches(0.35),
                desc, font_size=11, color=LIGHT_GRAY)

# =============================================
# PHASE 3: SYSTEM INTEGRATION
# =============================================
x3 = phase_start_x + 2 * phase_gap
add_rounded_rect(slide, x3, card_y, phase_width, card_h, BG_CARD, ACCENT_ORANGE)

items_phase3 = [
    ('移动端开发', 'Android/iOS轻量化APP'),
    ('功能集成', '摄像头调用、识别展示、历史查询'),
    ('部署验证', '嵌入式设备部署，确保稳定性'),
]

for j, (title, desc) in enumerate(items_phase3):
    item_y = card_y + Inches(0.18) + j * Inches(0.65)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x3 + Inches(0.2), item_y + Pt(6), Pt(8), Pt(8))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_ORANGE
    dot.line.fill.background()
    add_textbox(slide, x3 + Inches(0.45), item_y - Pt(2), Inches(3), Inches(0.28),
                title, font_size=13, color=WHITE, bold=True)
    add_textbox(slide, x3 + Inches(0.45), item_y + Inches(0.24), Inches(3.2), Inches(0.35),
                desc, font_size=11, color=LIGHT_GRAY)

# =============================================
# BOTTOM KEY POINTS BAR
# =============================================
bar_y = Inches(5.0)
bar_h = Inches(1.8)
add_rounded_rect(slide, Inches(0.8), bar_y, Inches(11.7), bar_h, BG_CARD, None)

add_textbox(slide, Inches(1.1), bar_y + Inches(0.15), Inches(3), Inches(0.35),
            '核心目标', font_size=16, color=ACCENT_PURPLE, bold=True)

key_points = [
    ('精准识别', ['规避传统字符分割的', '累积误差传递']),
    ('实时响应', ['参数调优平衡准确率', '与推理速度']),
    ('工程实用', ['轻量化部署适配', '复杂工地场景']),
    ('系统稳定', ['嵌入式验证确保', '现场可靠性']),
]

kp_start_x = Inches(1.1)
kp_gap = Inches(2.8)

for i, (title, desc_lines) in enumerate(key_points):
    kx = kp_start_x + i * kp_gap
    ky = bar_y + Inches(0.55)
    add_circle_icon(slide, kx, ky, Inches(0.38), ACCENT_PURPLE, '\u2713', font_size=14)
    add_textbox(slide, kx + Inches(0.48), ky - Pt(3), Inches(2.2), Inches(0.28),
                title, font_size=13, color=WHITE, bold=True)
    add_multiline_textbox(slide, kx + Inches(0.48), ky + Inches(0.22), Inches(2.2), Inches(0.7),
                          desc_lines, font_size=10, color=LIGHT_GRAY)

# Flow arrows
flow_y = card_y + card_h + Inches(0.12)
for i in range(2):
    fx = phase_start_x + phase_width + i * phase_gap + Inches(0.1)
    add_arrow_right(slide, fx, flow_y, Inches(0.28), Inches(0.14), RGBColor(0x3A, 0x3F, 0x5C))

# =============================================
# SAVE
# =============================================
output_path = r'C:\Users\HUAWEI\Desktop\FinalDesign\Python_VLPR-master\file\1.3.2_技术路线.pptx'
prs.save(output_path)
print(f'Saved: {output_path}')
